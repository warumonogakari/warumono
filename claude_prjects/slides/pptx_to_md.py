"""
pptx_to_md.py — PPTX → MD 変換スクリプト

使い方:
    uv run python slides/pptx_to_md.py <pptxファイルパス> [出力mdパス]

出力: 出力パス省略時はPPTXと同じフォルダに同名.mdを生成

出力フォーマット:
    # プレゼンタイトル（ファイル名）
    ---
    ## スライドN：タイトル
    #### テキスト（見出し扱い）
    - 箇条書き
    通常テキスト
    | col1 | col2 |  テーブル
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
try:
    from pptx.util import Pt
except ImportError:
    pass


def extract_table(shape):
    """テーブルをmarkdown形式に変換"""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = [cell.text_frame.text.replace('\n', ' ').strip() for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ''

    lines = []
    # ヘッダー行
    lines.append('| ' + ' | '.join(rows[0]) + ' |')
    lines.append('| ' + ' | '.join(['---'] * len(rows[0])) + ' |')
    for row in rows[1:]:
        lines.append('| ' + ' | '.join(row) + ' |')
    return '\n'.join(lines)


def extract_text_frame(shape):
    """テキストフレームからmarkdown行リストを返す"""
    if not shape.has_text_frame:
        return []

    lines = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # 箇条書き判定（インデントレベル > 0、またはbuChar設定あり）
        level = para.level
        pPr = para._p.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        has_bullet = False
        if pPr is not None:
            buChar = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
            buFont = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
            buAutoNum = pPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
            has_bullet = buChar is not None or buFont is not None or buAutoNum is not None

        indent = '  ' * level if level > 0 else ''
        if has_bullet or level > 0:
            lines.append(f'{indent}- {text}')
        else:
            lines.append(text)

    return lines


def is_title_shape(shape, slide):
    """タイトルプレースホルダーかどうかを判定"""
    try:
        from pptx.util import Pt
        from pptx.enum.text import PP_PLACEHOLDER
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            # TITLE=1, CENTER_TITLE=3
            if ph_type in (1, 3):
                return True
    except Exception:
        pass
    return False


def extract_slide(slide, slide_num):
    """1スライドをmarkdownブロックに変換"""
    title = ''
    content_blocks = []

    # タイトルシェイプを先に探す
    title_shape = None
    for shape in slide.shapes:
        if is_title_shape(shape, slide):
            title_shape = shape
            title = shape.text_frame.text.strip().replace('\n', ' ')
            break

    # タイトルが見つからない場合は最初のテキストボックスをタイトルとみなす
    if not title:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.strip().splitlines()[0]
                title_shape = shape
                break

    # コンテンツを上から順に処理（y座標でソート）
    shapes_sorted = sorted(
        [s for s in slide.shapes if s != title_shape],
        key=lambda s: (s.top, s.left)
    )

    for shape in shapes_sorted:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl_md = extract_table(shape)
            if tbl_md:
                content_blocks.append(tbl_md)
        elif shape.has_text_frame:
            lines = extract_text_frame(shape)
            if lines:
                content_blocks.append('\n'.join(lines))

    # 組み立て
    header = f'## スライド{slide_num}：{title}' if title else f'## スライド{slide_num}'
    parts = [header]
    for block in content_blocks:
        parts.append(block)

    return '\n\n'.join(parts)


def convert(pptx_path, md_path):
    prs = Presentation(str(pptx_path))
    stem = pptx_path.stem

    blocks = [f'# {stem}']
    for i, slide in enumerate(prs.slides, 1):
        blocks.append(extract_slide(slide, i))

    md_text = '\n\n---\n\n'.join(blocks) + '\n'
    md_path.write_text(md_text, encoding='utf-8')
    print(f'変換完了: {md_path}（{len(prs.slides)}枚）')


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('使い方: uv run python slides/pptx_to_md.py <pptxファイル> [出力mdパス]')
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    md_path = Path(sys.argv[2]) if len(sys.argv) >= 3 else pptx_path.with_suffix('.md')

    print(f'PPTX: {pptx_path}')
    convert(pptx_path, md_path)
