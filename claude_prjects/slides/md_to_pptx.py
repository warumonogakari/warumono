"""
md_to_pptx.py — MD → PPTX 自動変換スクリプト

使い方:
    uv run python slides/md_to_pptx.py <mdファイルパス> [出力pptxパス]

出力: 出力パス省略時はMDと同じフォルダに同名.pptxを生成

MDフォーマット規則:
    ## スライドN：タイトル     → スライドタイトル
    ### テキスト               → スライド末尾ならadd_quote_box、途中ならcb.h3()
    #### テキスト              → cb.h4()
    - 箇条書き                 → cb.bullet()（インデント続行行は結合）
    > 引用                     → cb.body()
    | col1 | col2 |            → cb.pair()（ヘッダー行はcb.h4()、セパレーターはスキップ）
    **写真...** **地図...**    → スキップ
    ※ ノート：...             → スキップ
    通常テキスト               → cb.body()
"""
import sys, re, copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ============================================================
# 定数
# ============================================================
GREEN       = RGBColor(0x32, 0xB4, 0x90)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE_S100 = RGBColor(0xFF, 0xA5, 0x00)  # ヘッダー・タイトル背景
ORANGE_ODD  = RGBColor(0xFF, 0xED, 0xCC)  # データ奇数行
ORANGE_EVEN = RGBColor(0xFF, 0xDB, 0x99)  # データ偶数行
TITLE_NAME  = "Google Shape;208;p31"
SLIDE_H     = 6858000

CONT_LEFT   = Emu(958715)
CONT_W      = Emu(10384788)

QUOTE_H     = Emu(900000)
QUOTE_BOT   = Emu(SLIDE_H - 150000)
QUOTE_TOP   = QUOTE_BOT - QUOTE_H
QUOTE_LEFT  = CONT_LEFT
QUOTE_W     = CONT_W

CONT_TOP_NORMAL        = Emu(1800000)
CONT_BOT_WITH_QUOTE    = QUOTE_TOP - Emu(100000)
CONT_BOT_WITHOUT_QUOTE = QUOTE_BOT

PT_TITLE = Pt(40)  # スライドタイトル（テンプレート既定44ptを上書き）
PT_H3,  PT_H4,  PT_BODY  = Pt(36), Pt(28), Pt(28)
PT_H3_C, PT_H4_C, PT_BODY_C = Pt(30), Pt(22), Pt(20)
PT_FOOT, PT_FOOT_C = Pt(16), Pt(14)  # 脚注（※始まりの箇条書き）
SP_SEC, SP_ITEM = Pt(14), Pt(4)

# TEMPLATE: 各スライドに流し込むデザイン（プレースホルダー配置）の供給元
TEMPLATE_PATH = "/Users/katouhiroshi/warumono/claude_prjects/slides/template_slide.pptx"
# BASE: スライドサイズ・テーマ・マスターを提供する「器」。中身スライドは init_prs() で全削除される。
# 成果物デッキとは独立した専用ファイルにすること（デッキ移動で道連れにならないように）。
BASE_PATH     = "/Users/katouhiroshi/warumono/claude_prjects/slides/base_theme.pptx"

# スキップキーワード（**テキスト** 独立行）
SKIP_KEYWORDS = ['写真', '地図', 'Photo', 'photo']


# ============================================================
# MDパーサー
# ============================================================

def strip_bold(text):
    return re.sub(r'\*\*(.*?)\*\*', r'\1', text)

def is_table_separator(line):
    return bool(re.match(r'^\|[-|:\s]+\|$', line.strip()))

def parse_table_cells(line):
    cells = [strip_bold(c.strip()) for c in line.strip().strip('|').split('|')]
    return [c for c in cells]  # 空文字も保持（列数判定のため）

def parse_slide_block(block):
    """1スライド分のブロックを解析してdictを返す"""
    lines = block.split('\n')
    title = ''
    items = []
    i = 0

    while i < len(lines):
        line = lines[i]

        # スライドタイトル: ## スライドN：タイトル / ## スライドN-M：タイトル / ## タイトル
        # 通番（スライド1-0：等、ハイフン入りも含む）は生成時に剥がす
        m = re.match(r'^## (?:スライド[\d\-‐－]+[：:]\s*)?(.+)$', line)
        if m:
            title = m.group(1).strip()
            i += 1
            continue

        # スキップ: ※ ノート：
        if line.strip().startswith('※ ノート'):
            i += 1
            continue

        # 空行スキップ
        if not line.strip():
            i += 1
            continue

        # h3: ### テキスト
        if line.startswith('### '):
            items.append({'type': 'h3', 'text': line[4:].strip()})
            i += 1
            continue

        # h4: #### テキスト
        if line.startswith('#### '):
            items.append({'type': 'h4', 'text': line[5:].strip()})
            i += 1
            continue

        # 箇条書き: - テキスト（次行がインデントなら結合）
        if line.startswith('- '):
            text = line[2:].strip()
            while i + 1 < len(lines) and re.match(r'^  +\S', lines[i + 1]):
                i += 1
                text += ' ' + lines[i].strip()
            items.append({'type': 'bullet', 'text': text})
            i += 1
            continue

        # 引用: > テキスト
        if line.startswith('> '):
            text = line[2:].strip()
            if text:
                items.append({'type': 'body', 'text': text})
            i += 1
            continue
        if line.strip() == '>':
            i += 1
            continue

        # テーブル行：ブロック全体を1アイテムとして収集
        if line.strip().startswith('|'):
            header = None
            rows = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                if is_table_separator(lines[i]):
                    i += 1
                    continue
                cells = [c.strip() for c in parse_table_cells(lines[i])]
                next_line = lines[i + 1] if i + 1 < len(lines) else ''
                if is_table_separator(next_line):
                    header = cells
                else:
                    rows.append(cells)
                i += 1
            if header is not None or rows:
                items.append({'type': 'table', 'header': header, 'rows': rows})
            continue

        # 太字のみの独立行: **テキスト**
        if re.match(r'^\*\*.+\*\*$', line.strip()):
            text = strip_bold(line.strip())
            if any(kw in text for kw in SKIP_KEYWORDS):
                i += 1
                continue
            items.append({'type': 'h4', 'text': text})
            i += 1
            continue

        # 通常テキスト
        if line.strip():
            items.append({'type': 'body', 'text': line.strip()})
        i += 1

    return {'title': title, 'items': items} if title else None


def determine_params(items):
    """
    has_quote / quote_text / compact を決定。
    末尾の h3 → add_quote_box に昇格。
    compact: h4 + bullet + pair の合計が8行超。
    """
    has_quote, quote_text = False, ''
    if items and items[-1]['type'] == 'h3':
        has_quote = True
        quote_text = items[-1]['text']
        items = items[:-1]

    count = sum(1 for it in items if it['type'] in ('h4', 'bullet', 'pair', 'body'))
    compact = count > 8

    return items, has_quote, quote_text, compact


def parse_md(md_path):
    text = Path(md_path).read_text(encoding='utf-8')
    # 先頭の # 行（プレゼンタイトル）を除去
    text = re.sub(r'^# .+\n', '', text)
    blocks = re.split(r'\n---\n', text)
    slides = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        slide = parse_slide_block(block)
        if slide and slide['title']:
            slides.append(slide)
    return slides


# ============================================================
# pptx ヘルパー
# ============================================================

def set_slide_bg(slide, rgb_color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


def set_title(slide, text, color=None):
    if color is None:
        color = GREEN
    for shape in slide.shapes:
        if shape.name == TITLE_NAME and shape.has_text_frame:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ''
            para = tf.paragraphs[0]
            if para.runs:
                run = para.runs[0]
            else:
                run = para.add_run()
            run.text = text
            run.font.color.rgb = color
            run.font.size = PT_TITLE
            return


def add_quote_box(slide, text, compact=False, color=None):
    if color is None:
        color = GREEN
    fs = PT_H3_C if compact else PT_H3
    box = slide.shapes.add_textbox(QUOTE_LEFT, QUOTE_TOP, QUOTE_W, QUOTE_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = fs
    r.font.bold = False
    r.font.color.rgb = color
    return box


TABLE_PRE_H  = Emu(800000)   # テーブル前コンテンツ領域の高さ
TABLE_ROW_H  = Emu(560000)   # テーブル各行の高さ


def add_pptx_table(slide, item, top, compact=False):
    """MDのtableアイテムをpptxのtableとして追加し、テーブル下端のEMU値を返す"""
    header = item.get('header') or []
    rows   = item.get('rows', [])
    n_cols = max(len(header), max((len(r) for r in rows), default=0))
    if n_cols == 0:
        return top
    n_rows = len(rows) + (1 if header else 0)
    table_h = TABLE_ROW_H * n_rows
    col_w   = CONT_W // n_cols
    fs = Pt(20) if compact else Pt(22)

    tbl = slide.shapes.add_table(n_rows, n_cols, CONT_LEFT, top, CONT_W, table_h).table

    # 列幅を均等に設定
    for ci in range(n_cols):
        tbl.columns[ci].width = col_w

    r_idx = 0
    if header:
        for ci, text in enumerate(header[:n_cols]):
            cell = tbl.cell(r_idx, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ORANGE_S100
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
            run.font.size = fs
            run.font.bold = False
            run.font.color.rgb = GREEN
        r_idx += 1

    for data_row_idx, row in enumerate(rows):
        row_bg = ORANGE_ODD if data_row_idx % 2 == 0 else ORANGE_EVEN
        for ci, text in enumerate(row[:n_cols]):
            cell = tbl.cell(r_idx, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_bg
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
            run.font.size = fs
            run.font.bold = False
            run.font.color.rgb = GREEN
        r_idx += 1

    # 枠線を GREEN に設定
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in [qn('a:lnL'), qn('a:lnR'), qn('a:lnT'), qn('a:lnB')]:
                ln = etree.SubElement(tcPr, tag)
                ln.set('w', '12700')
                solidFill = etree.SubElement(ln, qn('a:solidFill'))
                etree.SubElement(solidFill, qn('a:srgbClr')).set('val', '32B490')

    return top + table_h + Emu(150000)


class ContentBuilder:
    def __init__(self, slide, has_quote=False, compact=False, text_color=None):
        self.color = text_color if text_color is not None else GREEN
        bot = CONT_BOT_WITH_QUOTE if has_quote else CONT_BOT_WITHOUT_QUOTE
        self.box = slide.shapes.add_textbox(
            CONT_LEFT, CONT_TOP_NORMAL, CONT_W, bot - CONT_TOP_NORMAL
        )
        self.tf = self.box.text_frame
        self.tf.word_wrap = True
        self._first = True
        self.fs_h3   = PT_H3_C   if compact else PT_H3
        self.fs_h4   = PT_H4_C   if compact else PT_H4
        self.fs_body = PT_BODY_C  if compact else PT_BODY
        self.fs_foot = PT_FOOT_C  if compact else PT_FOOT

    def _new_para(self):
        if self._first:
            self._first = False
            return self.tf.paragraphs[0]
        return self.tf.add_paragraph()

    def h3(self, text, center=False):
        p = self._new_para()
        p.space_before = Emu(int(SP_SEC))
        if center:
            p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = self.fs_h3
        r.font.bold = False
        r.font.color.rgb = self.color

    def h4(self, text):
        p = self._new_para()
        p.space_before = Emu(int(SP_SEC))
        r = p.add_run()
        r.text = text
        r.font.size = self.fs_h4
        r.font.bold = False
        r.font.color.rgb = self.color

    def bullet(self, text):
        # ※始まりの行は脚注扱い：小フォント・行頭の「・」なし（※自体が目印）
        is_foot = text.startswith('※')
        p = self._new_para()
        p.space_before = Emu(int(SP_ITEM))
        if not is_foot:
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', '457200')
            pPr.set('indent', '-457200')
            buClr = etree.SubElement(pPr, qn('a:buClr'))
            clr_hex = f'{self.color[0]:02X}{self.color[1]:02X}{self.color[2]:02X}'
            etree.SubElement(buClr, qn('a:srgbClr')).set('val', clr_hex)
            etree.SubElement(pPr, qn('a:buChar')).set('char', '・')
        r = p.add_run()
        r.text = text
        r.font.size = self.fs_foot if is_foot else self.fs_body
        r.font.color.rgb = self.color

    def body(self, text, bold=False):
        p = self._new_para()
        p.space_before = Emu(int(SP_ITEM))
        r = p.add_run()
        r.text = text
        r.font.size = self.fs_body
        r.font.bold = bold
        r.font.color.rgb = self.color

    def pair(self, left, right):
        p = self._new_para()
        p.space_before = Emu(int(SP_SEC))
        rl = p.add_run()
        rl.text = left + '  →  '
        rl.font.size = self.fs_body
        rl.font.bold = True
        rl.font.color.rgb = self.color
        rr = p.add_run()
        rr.text = right
        rr.font.size = self.fs_body
        rr.font.color.rgb = self.color

    def render(self, item):
        t = item['type']
        if   t == 'h3':     self.h3(item['text'])
        elif t == 'h4':     self.h4(item['text'])
        elif t == 'bullet': self.bullet(item['text'])
        elif t == 'body':   self.body(item['text'])
        elif t == 'pair':   self.pair(item['left'], item['right'])


# ============================================================
# 重なりチェック・是正
# ============================================================

def check_and_fix_overlaps(prs, output_path):
    H = prs.slide_height.emu
    issues = []
    for si, slide in enumerate(prs.slides):
        shapes = sorted(
            [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()],
            key=lambda s: s.top
        )
        for i in range(len(shapes) - 1):
            s1, s2 = shapes[i], shapes[i + 1]
            b1 = s1.top + s1.height
            if b1 > s2.top:
                ov = b1 - s2.top
                issues.append(f'  スライド{si+1}: [{s1.name}]↔[{s2.name}] {ov//914}pt 重なり → 是正')
                s1.height = max(s1.height - ov - Emu(50000), Emu(500000))
        for s in shapes:
            if s.top + s.height > H:
                ov = s.top + s.height - H
                issues.append(f'  スライド{si+1}: [{s.name}] 下端 {ov//914}pt はみ出し → 是正')
                s.height = s.height - ov - Emu(50000)
    if issues:
        print('\n【重なり是正】')
        for msg in issues:
            print(msg)
        prs.save(str(output_path))
        print('是正後に再保存しました。')
    else:
        print('\n【重なりチェック】問題なし')


# ============================================================
# prs 初期化 & スライド複製
# ============================================================

def init_prs():
    tpl = Presentation(TEMPLATE_PATH)
    tpl_xml = copy.deepcopy(tpl.slides[0]._element)

    prs = Presentation(BASE_PATH)
    sldIdLst = prs.slides._sldIdLst
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = sldIdLst[i].get(qn('r:id'))
        prs.part.drop_rel(rId)
        del sldIdLst[i]

    def new_slide():
        s = prs.slides.add_slide(prs.slide_layouts[0])
        src = dst = None
        for e in tpl_xml.iter():
            if e.tag.endswith('}spTree'): src = e; break
        for e in s._element.iter():
            if e.tag.endswith('}spTree'): dst = e; break
        if src is not None and dst is not None:
            for c in list(dst): dst.remove(c)
            for c in src: dst.append(copy.deepcopy(c))
        slide = prs.slides[-1]
        for sh in [sh for sh in slide.shapes if sh.name != TITLE_NAME]:
            sh._element.getparent().remove(sh._element)
        return prs.slides[-1]

    return prs, new_slide


# ============================================================
# メイン変換
# ============================================================

def build_pptx(slides, output_path):
    prs, new_slide = init_prs()

    for data in slides:
        title = data['title']
        items = list(data['items'])

        # "タイトル" ラベルのスライドは最初のbody行をタイトルとして使う
        is_cover = title == 'タイトル'
        if is_cover and items and items[0]['type'] == 'body':
            title = items.pop(0)['text']

        items, has_quote, quote_text, compact = determine_params(items)

        s = new_slide()
        text_color = WHITE if is_cover else GREEN
        if is_cover:
            set_slide_bg(s, ORANGE_S100)
        set_title(s, title, color=text_color)

        has_table = any(it['type'] == 'table' for it in items)

        if has_table:
            # テーブル前後のアイテムを分割
            tbl_idx = next(i for i, it in enumerate(items) if it['type'] == 'table')
            pre  = items[:tbl_idx]
            tbl_item = items[tbl_idx]
            post = items[tbl_idx + 1:]

            # テーブル前コンテンツ
            if pre:
                cb = ContentBuilder(s, has_quote=False, compact=compact, text_color=text_color)
                cb.box.height = TABLE_PRE_H
                for it in pre:
                    cb.render(it)

            # テーブル本体
            tbl_top = CONT_TOP_NORMAL + (TABLE_PRE_H + Emu(100000) if pre else Emu(0))
            next_top = add_pptx_table(s, tbl_item, tbl_top, compact=compact)

            # テーブル後コンテンツ（quote_box は別途）
            if post:
                post_h = (QUOTE_TOP if has_quote else QUOTE_BOT) - next_top - Emu(50000)
                cb2 = ContentBuilder(s, has_quote=has_quote, compact=compact, text_color=text_color)
                cb2.box.top    = next_top
                cb2.box.height = max(post_h, Emu(300000))
                for it in post:
                    cb2.render(it)
        elif items:
            cb = ContentBuilder(s, has_quote=has_quote, compact=compact, text_color=text_color)
            for item in items:
                cb.render(item)

        if has_quote:
            add_quote_box(s, quote_text, compact=compact, color=text_color)

    prs.save(str(output_path))
    print(f'保存完了: {output_path}（{len(prs.slides)}枚）')
    check_and_fix_overlaps(prs, output_path)


# ============================================================
# エントリーポイント
# ============================================================

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('使い方: uv run python slides/md_to_pptx.py <mdファイル> [出力pptxパス]')
        sys.exit(1)

    md_path = Path(sys.argv[1])
    output_path = Path(sys.argv[2]) if len(sys.argv) >= 3 else md_path.with_suffix('.pptx')

    print(f'MD: {md_path}')
    slides = parse_md(md_path)
    print(f'スライド数: {len(slides)}枚')
    build_pptx(slides, output_path)
